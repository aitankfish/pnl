#!/usr/bin/env python3
"""
PNL Competition Presentation - FishTank Analogy
4-minute pitch deck with Business Model Canvas
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Theme Colors ───
DEEP_BLUE    = RGBColor(0x0A, 0x1A, 0x3A)  # Dark ocean
OCEAN_BLUE   = RGBColor(0x0D, 0x47, 0xA1)  # Primary blue
LIGHT_BLUE   = RGBColor(0x42, 0xA5, 0xF5)  # Accent
CYAN         = RGBColor(0x00, 0xE5, 0xFF)   # Neon accent
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xE0, 0xE0, 0xE0)
SOFT_WHITE   = RGBColor(0xF5, 0xF5, 0xF5)
GREEN        = RGBColor(0x00, 0xC8, 0x53)   # Support/YES
RED          = RGBColor(0xFF, 0x17, 0x44)   # Reject/NO
GOLD         = RGBColor(0xFF, 0xD6, 0x00)   # Revenue/fees
DARK_TEXT     = RGBColor(0x1A, 0x1A, 0x2E)
TEAL         = RGBColor(0x00, 0x96, 0x88)
SAND         = RGBColor(0xF0, 0xE6, 0xD3)   # River/sand

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None, border_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width or 1)
    else:
        shape.line.fill.background()
    return shape

def add_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width or 1)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_circle(slide, left, top, size, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    return shape

def set_shape_text(shape, text, font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = alignment
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)

def add_multi_text(slide, left, top, width, height, lines, default_size=16, default_color=WHITE):
    """lines = list of (text, size, color, bold, alignment)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        text = line_data[0]
        size = line_data[1] if len(line_data) > 1 else default_size
        color = line_data[2] if len(line_data) > 2 else default_color
        bold = line_data[3] if len(line_data) > 3 else False
        align = line_data[4] if len(line_data) > 4 else PP_ALIGN.LEFT
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Calibri"
        p.alignment = align
        p.space_before = Pt(4)
        p.space_after = Pt(4)
    return txBox

def add_arrow_shape(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


# ════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, DEEP_BLUE)

# Decorative water wave bar at bottom
add_rect(slide, Inches(0), Inches(6.5), Inches(13.333), Inches(1), OCEAN_BLUE)
add_rect(slide, Inches(0), Inches(6.4), Inches(13.333), Inches(0.15), LIGHT_BLUE)

# Fish emoji representations using shapes
# Small decorative fish (circles) scattered
positions = [(1, 2), (2.5, 5.5), (10, 1.5), (11.5, 4), (0.5, 4.5), (12, 6)]
for x, y in positions:
    c = add_circle(slide, Inches(x), Inches(y), Inches(0.3), LIGHT_BLUE)
    c.fill.fore_color.rgb = LIGHT_BLUE
    # Make semi-transparent by adjusting alpha isn't directly supported, so we use lighter colors

# Title text
add_text_box(slide, Inches(1.5), Inches(1.2), Inches(10), Inches(1.5),
             "PNL", 72, CYAN, True, PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1),
             "Predict & Launch", 36, WHITE, False, PP_ALIGN.CENTER)

# Tagline
add_text_box(slide, Inches(2), Inches(3.8), Inches(9), Inches(0.8),
             "The FishTank Where Ideas Swim or Sink", 28, LIGHT_BLUE, True, PP_ALIGN.CENTER)

# Subtitle
add_multi_text(slide, Inches(2.5), Inches(5), Inches(8), Inches(1.2), [
    ("Solana-Powered Prediction Market for Token Launches", 20, LIGHT_GRAY, False, PP_ALIGN.CENTER),
    ("", 10, WHITE),
    ("A decentralized platform where the community decides which projects deserve to launch", 16, RGBColor(0x90, 0xCA, 0xF9), False, PP_ALIGN.CENTER),
])


# ════════════════════════════════════════════════════════════════
# SLIDE 2: THE STORY — A Fish in a River
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DEEP_BLUE)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
             "The Story: A Fish With a Dream", 36, CYAN, True, PP_ALIGN.CENTER)

# River scene (left side)
add_rect(slide, Inches(0.3), Inches(1.5), Inches(3.8), Inches(5.2), RGBColor(0x1B, 0x5E, 0x20))  # riverbank
river_bg = add_shape(slide, Inches(0.5), Inches(1.7), Inches(3.5), Inches(4.8), RGBColor(0x0D, 0x47, 0x7A), LIGHT_BLUE, 2)

add_text_box(slide, Inches(0.7), Inches(1.9), Inches(3), Inches(0.5),
             "THE RIVER", 20, LIGHT_BLUE, True, PP_ALIGN.CENTER)

# Fish in river with idea
fish_circle = add_circle(slide, Inches(1.5), Inches(3.2), Inches(0.8), GOLD)
set_shape_text(fish_circle, "FISH", 12, DARK_TEXT, True)

# Thought bubble
thought = add_shape(slide, Inches(2.5), Inches(2.5), Inches(1.3), Inches(0.8), WHITE, OCEAN_BLUE, 1)
set_shape_text(thought, "I have\nan idea!", 10, DARK_TEXT, True)

add_multi_text(slide, Inches(0.6), Inches(4.5), Inches(3.2), Inches(2.5), [
    ("A small fish in a quiet river", 14, WHITE, True),
    ("has a project idea...", 14, LIGHT_GRAY),
    ("", 8, WHITE),
    ("But the river is small.", 13, RGBColor(0x90, 0xCA, 0xF9)),
    ("No audience. No funding.", 13, RGBColor(0x90, 0xCA, 0xF9)),
    ("No way to prove the idea.", 13, RGBColor(0x90, 0xCA, 0xF9)),
])

# Arrow
add_arrow_shape(slide, Inches(4.5), Inches(3.3), Inches(1.2), Inches(0.6), CYAN)

# FishTank scene (center)
tank_bg = add_shape(slide, Inches(5.9), Inches(1.5), Inches(3.8), Inches(5.2), RGBColor(0x0A, 0x27, 0x50), CYAN, 3)

add_text_box(slide, Inches(6.1), Inches(1.7), Inches(3.4), Inches(0.5),
             "THE FISHTANK (PNL)", 18, CYAN, True, PP_ALIGN.CENTER)

# Multiple fish inside tank - supporters (green) and predators (red)
supporter_positions = [(6.3, 2.8), (7.8, 3.8), (6.5, 4.5), (8.2, 2.6)]
for x, y in supporter_positions:
    c = add_circle(slide, Inches(x), Inches(y), Inches(0.5), GREEN)
    set_shape_text(c, "YES", 9, WHITE, True)

predator_positions = [(7.2, 2.4), (8.5, 4.2), (7.0, 4.8)]
for x, y in predator_positions:
    c = add_circle(slide, Inches(x), Inches(y), Inches(0.5), RED)
    set_shape_text(c, "NO", 9, WHITE, True)

# Our fish enters
our_fish = add_circle(slide, Inches(6.8), Inches(3.2), Inches(0.6), GOLD, CYAN)
set_shape_text(our_fish, "IDEA", 10, DARK_TEXT, True)

add_multi_text(slide, Inches(6.1), Inches(5.3), Inches(3.4), Inches(1.3), [
    ("Fish enters the tank!", 14, WHITE, True),
    ("Supporters & predators", 13, LIGHT_GRAY),
    ("vote on the idea", 13, LIGHT_GRAY),
])

# Arrow 2
add_arrow_shape(slide, Inches(10), Inches(3.3), Inches(1.2), Inches(0.6), GOLD)

# Ocean / Whale outcome (right)
ocean_bg = add_shape(slide, Inches(11.3), Inches(1.5), Inches(1.7), Inches(5.2), RGBColor(0x01, 0x57, 0x9B), GOLD, 2)

add_text_box(slide, Inches(11.3), Inches(1.7), Inches(1.7), Inches(0.5),
             "THE OCEAN", 14, GOLD, True, PP_ALIGN.CENTER)

# Whale
whale = add_circle(slide, Inches(11.5), Inches(3), Inches(1.2), GOLD)
set_shape_text(whale, "WHALE\n(Launched!)", 10, DARK_TEXT, True)

add_multi_text(slide, Inches(11.3), Inches(4.5), Inches(1.7), Inches(2), [
    ("If approved:", 12, GREEN, True),
    ("Fish becomes", 11, WHITE),
    ("a WHALE!", 14, GOLD, True),
    ("Token launches", 11, LIGHT_GRAY),
])


# ════════════════════════════════════════════════════════════════
# SLIDE 3: THE PROBLEM
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DEEP_BLUE)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
             "The Problem", 40, RED, True, PP_ALIGN.CENTER)

# Left column - For Creators
creator_box = add_shape(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.3), RGBColor(0x1A, 0x23, 0x40), RED, 2)
add_text_box(slide, Inches(0.8), Inches(1.7), Inches(5.2), Inches(0.6),
             "For Project Creators (The Fish)", 22, RED, True, PP_ALIGN.LEFT)

add_multi_text(slide, Inches(0.8), Inches(2.5), Inches(5.2), Inches(4), [
    ("No way to validate ideas before investing time & money", 18, WHITE, False),
    ("", 10, WHITE),
    ("Launching blind — 95% of token launches fail", 18, WHITE, False),
    ("", 10, WHITE),
    ("No community feedback before launch", 18, WHITE, False),
    ("", 10, WHITE),
    ("Rug pulls destroyed trust in new projects", 18, WHITE, False),
    ("", 10, WHITE),
    ("Result: Great ideas die in the river,", 18, GOLD, True),
    ("bad ones flood the ocean", 18, GOLD, True),
])

# Right column - For Community
community_box = add_shape(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.3), RGBColor(0x1A, 0x23, 0x40), LIGHT_BLUE, 2)
add_text_box(slide, Inches(7.1), Inches(1.7), Inches(5.2), Inches(0.6),
             "For Community (The School of Fish)", 22, LIGHT_BLUE, True, PP_ALIGN.LEFT)

add_multi_text(slide, Inches(7.1), Inches(2.5), Inches(5.2), Inches(4), [
    ("No voice in what gets launched", 18, WHITE),
    ("", 10, WHITE),
    ("Can't earn from early project discovery", 18, WHITE),
    ("", 10, WHITE),
    ("Scams & low-effort projects waste time", 18, WHITE),
    ("", 10, WHITE),
    ("No skin-in-the-game prediction mechanism", 18, WHITE),
    ("", 10, WHITE),
    ("Result: Community has no power", 18, GOLD, True),
    ("to filter signal from noise", 18, GOLD, True),
])


# ════════════════════════════════════════════════════════════════
# SLIDE 4: HOW IT WORKS — The FishTank Flow
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DEEP_BLUE)

add_text_box(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
             "How The FishTank Works", 36, CYAN, True, PP_ALIGN.CENTER)

# Step boxes
steps = [
    ("1", "CREATE", "Fish fills the form\n\nProject name, description,\ntoken symbol, category,\ntarget pool, team info,\nsocial links, pitch video", GREEN, "Creator submits\nproject to the tank"),
    ("2", "PREDICT", "School of fish votes\n\nYES = Support the idea\nNO = Reject it\n\nPut SOL behind\nyour prediction", LIGHT_BLUE, "Community stakes\nSOL on outcome"),
    ("3", "RESOLVE", "Timer expires\n\nMarket resolves based\non community vote\n\nWinning side earns\nfrom losing side", GOLD, "Winners collect,\nlosers learn"),
    ("4", "LAUNCH", "Fish becomes a whale!\n\nApproved projects auto-\nlaunch on pump.fun\n\nToken goes live\non Solana", CYAN, "Token launches\nautomatically"),
]

step_width = 2.8
gap = 0.25
start_x = 0.5

for i, (num, title, desc, color, subtitle) in enumerate(steps):
    x = start_x + i * (step_width + gap)

    # Step card
    card = add_shape(slide, Inches(x), Inches(1.2), Inches(step_width), Inches(5.5), RGBColor(0x0F, 0x1F, 0x3D), color, 2)

    # Step number circle
    num_circle = add_circle(slide, Inches(x + step_width/2 - 0.3), Inches(1.4), Inches(0.6), color)
    set_shape_text(num_circle, num, 24, WHITE if color != GOLD else DARK_TEXT, True)

    # Title
    add_text_box(slide, Inches(x + 0.1), Inches(2.1), Inches(step_width - 0.2), Inches(0.5),
                 title, 22, color, True, PP_ALIGN.CENTER)

    # Description
    add_multi_text(slide, Inches(x + 0.15), Inches(2.7), Inches(step_width - 0.3), Inches(3), [
        (desc, 13, WHITE),
    ])

    # Subtitle at bottom
    add_text_box(slide, Inches(x + 0.1), Inches(5.8), Inches(step_width - 0.2), Inches(0.8),
                 subtitle, 12, color, True, PP_ALIGN.CENTER)

    # Arrow between steps
    if i < 3:
        arrow_x = x + step_width
        add_arrow_shape(slide, Inches(arrow_x), Inches(3.5), Inches(gap + 0.05), Inches(0.3), color)

# Bottom note
add_text_box(slide, Inches(1), Inches(6.85), Inches(11), Inches(0.5),
             "[ Screenshots of actual PNL platform: Create Form  |  Market Page  |  Voting Interface  |  Launched Tokens ]",
             14, RGBColor(0x64, 0xB5, 0xF6), False, PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 5: DEMO — Platform Screenshots Placeholder
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DEEP_BLUE)

add_text_box(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
             "Live Platform Demo", 36, CYAN, True, PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.5), Inches(0.85), Inches(12), Inches(0.5),
             "pnl.fun — Built & Deployed on Solana Mainnet", 18, LIGHT_GRAY, False, PP_ALIGN.CENTER)

# Screenshot placeholders with labels
screenshots = [
    ("Create Project Form", "Fish fills out the form to\nenter the tank\n\n[ INSERT SCREENSHOT\nof /create page ]", Inches(0.4), Inches(1.5), Inches(4), Inches(5.3)),
    ("Market / FishTank", "Community votes YES or NO\non the project\n\n[ INSERT SCREENSHOT\nof /market/[id] page ]", Inches(4.6), Inches(1.5), Inches(4), Inches(5.3)),
    ("Launched Tokens", "Successful fish that became\nwhales on pump.fun\n\n[ INSERT SCREENSHOT\nof /launched page ]", Inches(8.8), Inches(1.5), Inches(4), Inches(5.3)),
]

for title, desc, left, top, width, height in screenshots:
    # Screenshot frame
    frame = add_shape(slide, left, top, width, height, RGBColor(0x15, 0x25, 0x45), CYAN, 2)

    # Title bar
    add_rect(slide, left + Inches(0.05), top + Inches(0.05), width - Inches(0.1), Inches(0.4), OCEAN_BLUE)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.08), width - Inches(0.3), Inches(0.35),
                 title, 14, WHITE, True, PP_ALIGN.CENTER)

    # Placeholder text
    add_multi_text(slide, left + Inches(0.2), top + Inches(0.8), width - Inches(0.4), height - Inches(1), [
        (desc, 14, LIGHT_GRAY, False, PP_ALIGN.CENTER),
        ("", 10, WHITE),
        ("Replace this with an actual", 12, RGBColor(0x64, 0xB5, 0xF6), False, PP_ALIGN.CENTER),
        ("screenshot from pnl.fun", 12, RGBColor(0x64, 0xB5, 0xF6), False, PP_ALIGN.CENTER),
    ])


# ════════════════════════════════════════════════════════════════
# SLIDE 6: BUSINESS MODEL CANVAS
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_text_box(slide, Inches(0.3), Inches(0.15), Inches(12.5), Inches(0.6),
             "Business Model Canvas — PNL (Predict & Launch)", 30, DARK_TEXT, True, PP_ALIGN.CENTER)

# Canvas grid - Strategyzer layout
# Row 1: Key Partners | Key Activities | Value Prop | Customer Rel | Customer Segments
# Row 1b:             | Key Resources  |            | Channels     |
# Row 2: Cost Structure                | Revenue Streams

border_c = RGBColor(0x90, 0x90, 0x90)
bw = 1.5
cell_bg = RGBColor(0xFA, 0xFA, 0xFF)

# Dimensions
canvas_top = Inches(0.85)
canvas_left = Inches(0.3)
col_w = [Inches(2.4), Inches(2.4), Inches(2.8), Inches(2.4), Inches(2.4)]
row_h_top = Inches(2.7)
row_h_bot = Inches(1.6)
total_w = sum(c for c in col_w)

# --- ROW 1 TOP HALF (Key Partners, Key Activities top, Value Prop, Cust Rel top, Cust Segments) ---

canvas_data_top = [
    ("Key Partners", "Solana Foundation\npump.fun (launch partner)\nJupiter (DEX aggregator)\nHelius (RPC/WebSocket)\nPinata (IPFS storage)\nPrivy (Auth)"),
    ("Key Activities", "Platform development\nSmart contract audits\nCommunity growth\nMarket curation\nLaunch pipeline mgmt"),
    ("Value Propositions", "Community-validated\ntoken launches\n\nEarn by predicting right\n\nFilter signal from noise\n\nDemocratized fundraising\nvia the FishTank"),
    ("Customer Relationships", "Real-time chat rooms\nVoice rooms per market\nLive activity feeds\nFollower/following system\nNotifications"),
    ("Customer Segments", "Crypto project creators\n(the Fish)\n\nDegen traders/predictors\n(the School)\n\nSolana ecosystem\nbuilders & investors"),
]

canvas_data_bottom_left = [
    ("Key Resources", "Solana smart contract\nWeb platform (Next.js)\nMongoDB + Redis\nSocket.IO real-time\nCommunity of predictors"),
]

canvas_data_bottom_right = [
    ("Channels", "pnl.fun (web app)\nTwitter/X\nTelegram community\nSolana ecosystem events\nPartner integrations"),
]

# Draw top row cells
x_pos = canvas_left
for i, (title, content) in enumerate(canvas_data_top):
    # For col 1 and 4 (Key Activities/Resources and Cust Rel/Channels), split vertically
    if i in [1, 3]:
        # Top half only
        cell = add_rect(slide, x_pos, canvas_top, col_w[i], Inches(1.35), cell_bg, border_c, bw)
    elif i == 2:
        # Value prop - full height
        cell = add_rect(slide, x_pos, canvas_top, col_w[i], row_h_top, cell_bg, border_c, bw)
    else:
        # Key Partners and Customer Segments - full height
        cell = add_rect(slide, x_pos, canvas_top, col_w[i], row_h_top, cell_bg, border_c, bw)

    # Title
    add_text_box(slide, x_pos + Inches(0.08), canvas_top + Inches(0.03), col_w[i] - Inches(0.15), Inches(0.3),
                 title, 10, OCEAN_BLUE, True, PP_ALIGN.LEFT)

    # Content
    content_top = canvas_top + Inches(0.3)
    if i in [1, 3]:
        content_h = Inches(0.95)
    else:
        content_h = row_h_top - Inches(0.35)

    add_text_box(slide, x_pos + Inches(0.08), content_top, col_w[i] - Inches(0.15), content_h,
                 content, 9, DARK_TEXT, False, PP_ALIGN.LEFT)

    x_pos += col_w[i]

# Draw bottom-left sub-cells (Key Resources under Key Activities)
x_kr = canvas_left + col_w[0]
kr_top = canvas_top + Inches(1.35)
cell = add_rect(slide, x_kr, kr_top, col_w[1], Inches(1.35), cell_bg, border_c, bw)
add_text_box(slide, x_kr + Inches(0.08), kr_top + Inches(0.03), col_w[1] - Inches(0.15), Inches(0.3),
             "Key Resources", 10, OCEAN_BLUE, True, PP_ALIGN.LEFT)
add_text_box(slide, x_kr + Inches(0.08), kr_top + Inches(0.3), col_w[1] - Inches(0.15), Inches(0.95),
             canvas_data_bottom_left[0][1], 9, DARK_TEXT, False, PP_ALIGN.LEFT)

# Draw bottom-right sub-cells (Channels under Customer Relationships)
x_ch = canvas_left + col_w[0] + col_w[1] + col_w[2]
cell = add_rect(slide, x_ch, kr_top, col_w[3], Inches(1.35), cell_bg, border_c, bw)
add_text_box(slide, x_ch + Inches(0.08), kr_top + Inches(0.03), col_w[3] - Inches(0.15), Inches(0.3),
             "Channels", 10, OCEAN_BLUE, True, PP_ALIGN.LEFT)
add_text_box(slide, x_ch + Inches(0.08), kr_top + Inches(0.3), col_w[3] - Inches(0.15), Inches(0.95),
             canvas_data_bottom_right[0][1], 9, DARK_TEXT, False, PP_ALIGN.LEFT)

# --- ROW 2: Cost Structure | Revenue Streams ---
row2_top = canvas_top + row_h_top + Inches(0.05)
cost_w = col_w[0] + col_w[1] + Inches(col_w[2] / 2)
rev_w = total_w - cost_w

# Cost Structure
cell = add_rect(slide, canvas_left, row2_top, cost_w, row_h_bot, cell_bg, border_c, bw)
add_text_box(slide, canvas_left + Inches(0.08), row2_top + Inches(0.03), cost_w - Inches(0.15), Inches(0.3),
             "Cost Structure", 10, OCEAN_BLUE, True, PP_ALIGN.LEFT)
add_text_box(slide, canvas_left + Inches(0.08), row2_top + Inches(0.3), cost_w - Inches(0.15), row_h_bot - Inches(0.35),
             "Cloud hosting (Render, Vercel)\nSolana RPC costs (Helius)\nSmart contract audits\nTeam salaries\nMarketing & community growth\nIPFS storage (Pinata)", 9, DARK_TEXT, False, PP_ALIGN.LEFT)

# Revenue Streams
cell = add_rect(slide, canvas_left + cost_w, row2_top, rev_w, row_h_bot, RGBColor(0xFF, 0xFD, 0xE7), border_c, bw)
add_text_box(slide, canvas_left + cost_w + Inches(0.08), row2_top + Inches(0.03), rev_w - Inches(0.15), Inches(0.3),
             "Revenue Streams  (The Fish Waste)", 10, RGBColor(0xE6, 0x5C, 0x00), True, PP_ALIGN.LEFT)
add_text_box(slide, canvas_left + cost_w + Inches(0.08), row2_top + Inches(0.3), rev_w - Inches(0.15), row_h_bot - Inches(0.35),
             "Market creation fees (SOL)\nTrading fees on predictions\nToken launch fees\nPlatform cut on market resolution\nMerch store revenue (SOL payments)\nPremium features (future)", 9, DARK_TEXT, False, PP_ALIGN.LEFT)

# Footer
add_text_box(slide, Inches(0.3), Inches(7), Inches(12), Inches(0.4),
             "PNL — Predict & Launch  |  pnl.fun  |  Built on Solana", 11, RGBColor(0x99, 0x99, 0x99), False, PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 7: REVENUE MODEL — The Fish Waste (Fees & Tokenomics)
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DEEP_BLUE)

add_text_box(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
             "Revenue Model: The Fish Waste", 36, GOLD, True, PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(0.9), Inches(11), Inches(0.5),
             "When fish are in the water, they excrete waste — that's the fees that power the ecosystem",
             16, LIGHT_GRAY, False, PP_ALIGN.CENTER)

# Revenue streams as visual blocks
rev_items = [
    ("Market Creation Fee", "Creators pay SOL to\nenter the FishTank\n\nFilters spam projects\nand funds the platform", "0.05 SOL", TEAL),
    ("Prediction Fees", "Small fee on every\nYES/NO prediction\n\nTaken from both sides\nwhen market resolves", "2% of stake", OCEAN_BLUE),
    ("Launch Fee", "When a fish becomes\na whale (token launches)\n\nPlatform takes a cut\nof the launch pool", "1% of pool", GREEN),
    ("Token Ownership", "Supporting fish can claim\ntokens from launched projects\n\nEarly backers rewarded\nwith real token allocation", "Pro-rata share", GOLD),
]

for i, (title, desc, amount, color) in enumerate(rev_items):
    x = Inches(0.4 + i * 3.2)

    # Card
    card = add_shape(slide, x, Inches(1.6), Inches(2.9), Inches(4.5), RGBColor(0x0F, 0x1F, 0x3D), color, 2)

    # Title
    add_text_box(slide, x + Inches(0.1), Inches(1.8), Inches(2.7), Inches(0.5),
                 title, 18, color, True, PP_ALIGN.CENTER)

    # Amount badge
    badge = add_shape(slide, x + Inches(0.6), Inches(2.4), Inches(1.7), Inches(0.5), color)
    set_shape_text(badge, amount, 16, WHITE if color != GOLD else DARK_TEXT, True)

    # Description
    add_text_box(slide, x + Inches(0.15), Inches(3.1), Inches(2.6), Inches(2.8),
                 desc, 13, WHITE, False, PP_ALIGN.CENTER)

# Bottom summary
summary_box = add_shape(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.9), RGBColor(0x1A, 0x23, 0x40), GOLD, 2)
add_multi_text(slide, Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.7), [
    ("Revenue flows back to: Platform Treasury (operational costs) + Project Founders (launch funding) + Community (token rewards)", 14, GOLD, True, PP_ALIGN.CENTER),
])


# ════════════════════════════════════════════════════════════════
# SLIDE 8: TRACTION & ASK
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DEEP_BLUE)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
             "Why PNL Wins", 40, CYAN, True, PP_ALIGN.CENTER)

# Left: What we've built
built_box = add_shape(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(3), RGBColor(0x0F, 0x1F, 0x3D), GREEN, 2)
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5.4), Inches(0.5),
             "What We've Built (Live on Mainnet)", 20, GREEN, True, PP_ALIGN.LEFT)

add_multi_text(slide, Inches(0.8), Inches(2.1), Inches(5.4), Inches(2), [
    ("Deployed Solana smart contract (mainnet)", 15, WHITE),
    ("Full-stack web platform at pnl.fun", 15, WHITE),
    ("Real-time chat, voice rooms, live feeds", 15, WHITE),
    ("Jupiter swap integration", 15, WHITE),
    ("pump.fun auto-launch pipeline", 15, WHITE),
    ("Merch store with SOL payments", 15, WHITE),
])

# Right: Why now
why_box = add_shape(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(3), RGBColor(0x0F, 0x1F, 0x3D), LIGHT_BLUE, 2)
add_text_box(slide, Inches(7.1), Inches(1.5), Inches(5.4), Inches(0.5),
             "Why Now?", 20, LIGHT_BLUE, True, PP_ALIGN.LEFT)

add_multi_text(slide, Inches(7.1), Inches(2.1), Inches(5.4), Inches(2), [
    ("Solana is the #1 chain for launches", 15, WHITE),
    ("pump.fun proved the demand", 15, WHITE),
    ("But no curation layer exists yet", 15, WHITE),
    ("Community wants voice in what launches", 15, WHITE),
    ("Prediction markets are proven (Polymarket)", 15, WHITE),
    ("PNL = Polymarket meets pump.fun", 15, WHITE),
])

# Bottom CTA
cta_box = add_shape(slide, Inches(2), Inches(4.8), Inches(9.3), Inches(2.3), RGBColor(0x0A, 0x27, 0x50), CYAN, 3)

add_text_box(slide, Inches(2.3), Inches(5), Inches(8.7), Inches(0.6),
             "The FishTank is Open. Dive In.", 28, CYAN, True, PP_ALIGN.CENTER)

add_multi_text(slide, Inches(2.5), Inches(5.7), Inches(8.3), Inches(1.2), [
    ("pnl.fun  |  Built on Solana  |  Live on Mainnet", 20, WHITE, True, PP_ALIGN.CENTER),
    ("", 8, WHITE),
    ("Every great project deserves a community vote before it swims into the ocean.", 16, LIGHT_GRAY, False, PP_ALIGN.CENTER),
])

# Contact / Team
add_text_box(slide, Inches(1), Inches(7.1), Inches(11), Inches(0.3),
             "Thank you!  |  Questions?", 14, RGBColor(0x64, 0xB5, 0xF6), False, PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════
import os
REPO_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
output_path = os.path.join(REPO_ROOT, "docs", "pitch", "PNL_FishTank_Pitch.pptx")
os.makedirs(os.path.dirname(output_path), exist_ok=True)
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
